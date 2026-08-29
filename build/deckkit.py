"""deckkit - the FEU lecture deck toolkit for FIN1209.

This module knows how to draw slides. It knows nothing about any particular
chapter. Chapter content lives in its own module (see content_chapter01.py)
as plain data, so adding Chapter 2 is a new content file rather than a rewrite.

Design rules encoded here, from the course teaching brief:

  * One idea per slide, six lines of body text at the very most.
  * Every content slide carries a progress marker: "Part 3 of 6 - Classifications".
  * Green carries structure, gold marks the single thing to notice.
  * Check slides look nothing like teaching slides, so the room knows
    instantly that a question is coming.
  * Two or three speaker cue lines in the notes of every content slide.
  * A figure is a slide of its own: title, image, credit line. It renders the
    same with or without the artwork, so the deck built from the public
    repository keeps its slide count and its progress markers.
"""

from __future__ import annotations

import os
import zipfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# FEU institutional identity
# --------------------------------------------------------------------------

GREEN = RGBColor(0x00, 0x7A, 0x33)  # institution, primary, structure
GREEN_DEEP = RGBColor(0x00, 0x4F, 0x21)  # check slides, section dividers
GOLD = RGBColor(0xF2, 0xA9, 0x00)  # emphasis, achievement, the one thing to notice
INK = RGBColor(0x1A, 0x1A, 0x1A)  # body text
PAPER = RGBColor(0xFA, 0xF8, 0xF2)  # warm page ground
MUTED = RGBColor(0x81, 0x81, 0x81)  # captions, progress markers
BORDER = RGBColor(0xD7, 0xE0, 0xD9)  # soft rules and hairlines
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Display face for titles and section dividers. The FEU identity face is
# Marcellus SC, shared with AMS0011. It is not installed on the build machine
# or on the lecture room PCs, and PowerPoint substitutes a sans for it, which
# loses the display serif entirely. So the deck ships on the named serif
# fallback, Palatino. Once Marcellus SC is installed on the presenting
# machine, rebuild with:
#     .venv/bin/python build/build_chapter1.py --display-font "Marcellus SC"
DISPLAY_FONT = "Palatino"
BODY_FONT = "Arial"  # system sans, present on every teaching machine
MONO_FONT = "Courier New"  # answer keys and data

# Fixed timestamp so a rebuild of unchanged content is byte identical.
EPOCH = datetime(2020, 1, 1, 0, 0, 0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.85)
CONTENT_W = Emu(SLIDE_W - 2 * MARGIN)

# Where render_figure looks for artwork. The textbook figures are Wiley's and
# are not in this repository, so the folder is gitignored and is usually
# absent. When a figure file is missing the slide still renders, carrying the
# figure number, the credit line and the speaker cue, and the deck keeps its
# slide count and its progress markers. See chapter-01/README.md.
FIGURES_DIR: Path | None = None


# --------------------------------------------------------------------------
# Content model - what a chapter module hands to the builder
# --------------------------------------------------------------------------


@dataclass
class Slide:
    """Base for every slide the chapter content declares."""

    notes: tuple[str, ...] = field(default_factory=tuple)


@dataclass
class TitleSlide(Slide):
    course: str = ""
    chapter: str = ""
    title: str = ""
    subtitle: str = ""
    presenter: str = ""


@dataclass
class SectionOpen(Slide):
    number: int = 0
    total: int = 0
    title: str = ""
    covers: tuple[str, ...] = ()
    minutes: str = ""


@dataclass
class Content(Slide):
    """A teaching slide. Six body lines maximum, three preferred."""

    title: str = ""
    lines: tuple[str, ...] = ()
    accent: str = ""  # the one gold line, the thing to notice
    caption: str = ""  # small muted line under the accent


@dataclass
class Term(Slide):
    """A new term, taught plain words first and formal definition last."""

    term: str = ""
    plain: str = ""
    example: str = ""
    formal: str = ""


@dataclass
class Quote(Slide):
    """A short attributed quotation of a famous definition."""

    text: str = ""
    source: str = ""
    takeaway: str = ""


@dataclass
class Figure(Slide):
    """A textbook figure, one to a slide.

    ``number`` is the book's own figure number, "1.11". ``shows`` is one line
    saying what is on the image; it is what the placeholder prints when the
    artwork is absent, so it has to stand on its own as a description.
    """

    title: str = ""
    number: str = ""
    shows: str = ""

    @property
    def filename(self) -> str:
        major, minor = self.number.split(".")
        return f"figure-{major}-{int(minor):02d}.png"

    @property
    def credit(self) -> str:
        return (f"Figure {self.number} - Lim, The Handbook of Technical "
                "Analysis (Wiley, 2016)")


@dataclass
class Question:
    stem: str
    options: tuple[str, str, str, str]
    answer: str  # "A" / "B" / "C" / "D"
    reason: str


@dataclass
class Check(Slide):
    """Two multiple choice questions. Rendered as a question slide plus a
    reveal slide, so it always costs two slides in the deck."""

    label: str = ""
    questions: tuple[Question, ...] = ()


@dataclass
class Recap(Slide):
    """The you-now-know close of a section."""

    title: str = "You now know"
    items: tuple[str, ...] = ()


@dataclass
class Closing(Slide):
    title: str = ""
    lines: tuple[str, ...] = ()
    accent: str = ""


@dataclass
class Section:
    number: int
    title: str
    short: str  # the word used in the progress marker
    minutes: str
    covers: tuple[str, ...]
    slides: tuple[Slide, ...]
    recap: Recap


@dataclass
class Chapter:
    course: str
    code: str
    chapter: str
    title: str
    subtitle: str
    presenter: str
    objectives: tuple[str, ...]
    roadmap: tuple[str, ...]
    sections: tuple[Section, ...]
    closing: tuple[Slide, ...]


# --------------------------------------------------------------------------
# Drawing primitives
# --------------------------------------------------------------------------


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    # An autoshape carries a theme style reference that renders a drop shadow.
    # Strip it: the identity has no shadows, gradients or effects anywhere.
    style = shape._element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}style")
    if style is not None:
        shape._element.remove(style)
    shape.shadow.inherit = False
    return shape


# Rough text metrics, good enough to lay out a slide without a font engine.
# Average glyph advance as a fraction of the point size.
_ADVANCE = {"regular": 0.50, "bold": 0.53, "mono": 0.60}


def _line_count(text: str, width_in: float, size_pt: float,
                weight: str = "regular") -> int:
    """How many wrapped lines this text needs in a box of the given width."""
    if not text:
        return 0
    per_line = max(1.0, width_in * 72.0 / (size_pt * _ADVANCE[weight]))
    words, lines, current = text.split(), 1, 0
    for word in words:
        need = len(word) + (1 if current else 0)
        if current + need > per_line and current:
            lines += 1
            current = len(word)
        else:
            current += need
    return lines


def _text_height(text: str, width_in: float, size_pt: float,
                 weight: str = "regular", spacing: float = 1.15) -> float:
    """Height in inches of the wrapped text."""
    return _line_count(text, width_in, size_pt, weight) * size_pt * spacing / 72.0


def _ground(slide, color):
    """Paint the whole slide, since a blank layout is transparent."""
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)


def _textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    return frame


def _para(frame, text, *, size, color, font=BODY_FONT, bold=False,
          italic=False, align=PP_ALIGN.LEFT, space_before=0, space_after=0,
          first=False, line_spacing=1.0, hanging=None):
    para = frame.paragraphs[0] if first else frame.add_paragraph()
    para.alignment = align
    if hanging is not None:
        # Wrapped lines align under the text, not under the bullet mark.
        pPr = para._pPr if para._pPr is not None else para._p.get_or_add_pPr()
        pPr.set("marL", str(int(hanging)))
        pPr.set("indent", str(int(-hanging)))
    para.space_before = Pt(space_before)
    para.space_after = Pt(space_after)
    para.line_spacing = line_spacing
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return para


def _notes(slide, cues):
    """Speaker cues, so the instructor always knows the next sentence."""
    if not cues:
        return
    frame = slide.notes_slide.notes_text_frame
    frame.text = ""
    for i, cue in enumerate(cues):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = "- " + cue
        run.font.size = Pt(14)


def _progress(slide, marker, *, on_dark=False):
    """The small position marker every content slide carries."""
    frame = _textbox(slide, MARGIN, SLIDE_H - Inches(0.72), CONTENT_W, Inches(0.34))
    _para(
        frame,
        marker,
        size=12,
        color=WHITE if on_dark else MUTED,
        first=True,
    )


def _header_rule(slide, y):
    _rect(slide, MARGIN, y, CONTENT_W, Pt(1.25), BORDER)


def _gold_rule(slide, y, width=Inches(1.6)):
    _rect(slide, MARGIN, y, width, Pt(4), GOLD)


def _chip(slide, x, y, text, *, fill, text_color, width=Inches(1.9)):
    shape = _rect(slide, x, y, width, Inches(0.36), fill)
    frame = shape.text_frame
    frame.word_wrap = False
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(frame, text, size=12, color=text_color, bold=True,
          align=PP_ALIGN.CENTER, first=True)
    return shape


# --------------------------------------------------------------------------
# Slide renderers
# --------------------------------------------------------------------------


def render_title(prs, s: TitleSlide):
    slide = _blank(prs)
    _ground(slide, GREEN)
    _rect(slide, 0, SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.55), GREEN_DEEP)

    frame = _textbox(slide, MARGIN, Inches(1.55), CONTENT_W, Inches(0.5))
    _para(frame, s.course, size=15, color=GOLD, bold=True, first=True)

    _gold_rule(slide, Inches(2.25))

    frame = _textbox(slide, MARGIN, Inches(2.65), CONTENT_W, Inches(2.4))
    _para(frame, s.chapter, size=20, color=WHITE, font=DISPLAY_FONT, first=True)
    _para(frame, s.title, size=44, color=WHITE, font=DISPLAY_FONT,
          space_before=10, line_spacing=1.05)

    frame = _textbox(slide, MARGIN, Inches(5.35), CONTENT_W, Inches(1.0))
    _para(frame, s.subtitle, size=16, color=WHITE, first=True)
    _para(frame, s.presenter, size=14, color=GOLD, space_before=8)

    _notes(slide, s.notes)


def render_section_open(prs, s: SectionOpen):
    slide = _blank(prs)
    _ground(slide, GREEN)
    _rect(slide, 0, 0, Inches(0.28), SLIDE_H, GOLD)

    frame = _textbox(slide, MARGIN, Inches(1.15), CONTENT_W, Inches(0.4))
    _para(frame, f"Part {s.number} of {s.total}", size=15, color=GOLD,
          bold=True, first=True)

    frame = _textbox(slide, MARGIN, Inches(1.75), CONTENT_W, Inches(1.2))
    _para(frame, s.title, size=40, color=WHITE, font=DISPLAY_FONT,
          first=True, line_spacing=1.05)

    _rect(slide, MARGIN, Inches(3.2), Inches(1.6), Pt(3), GOLD)

    frame = _textbox(slide, MARGIN, Inches(3.7), CONTENT_W, Inches(2.2))
    for i, line in enumerate(s.covers):
        _para(frame, line, size=19, color=WHITE, first=(i == 0),
              space_after=10, line_spacing=1.15)

    frame = _textbox(slide, MARGIN, SLIDE_H - Inches(0.9), CONTENT_W, Inches(0.4))
    _para(frame, s.minutes, size=12, color=GOLD, first=True)

    _notes(slide, s.notes)


def _content_frame(slide, title, marker):
    """Shared chrome for every paper-ground teaching slide."""
    _ground(slide, PAPER)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.16), GREEN)

    frame = _textbox(slide, MARGIN, Inches(0.72), CONTENT_W, Inches(1.0))
    _para(frame, title, size=30, color=GREEN, font=DISPLAY_FONT, first=True,
          line_spacing=1.05)
    _header_rule(slide, Inches(1.72))
    _progress(slide, marker)


def render_content(prs, s: Content, marker):
    slide = _blank(prs)
    _content_frame(slide, s.title, marker)

    inner_in = CONTENT_W.inches
    top = 2.15
    if s.lines:
        frame = _textbox(slide, MARGIN, Inches(top), CONTENT_W, Inches(3.4))
        for i, line in enumerate(s.lines):
            _para(frame, line, size=21, color=INK, first=(i == 0),
                  space_after=13, line_spacing=1.18)
        top += sum(_text_height(line, inner_in, 21, spacing=1.18) + 0.20
                   for line in s.lines) + 0.34

    if s.accent:
        accent_h = _text_height(s.accent, inner_in - 0.32, 22, "bold") + 0.16
        _rect(slide, MARGIN, Inches(top), Inches(0.09), Inches(accent_h), GOLD)
        frame = _textbox(slide, MARGIN + Inches(0.32), Inches(top + 0.04),
                         CONTENT_W - Inches(0.32), Inches(accent_h))
        _para(frame, s.accent, size=22, color=INK, bold=True, first=True,
              line_spacing=1.15)
        top += accent_h + 0.26

    if s.caption:
        frame = _textbox(slide, MARGIN, Inches(top), CONTENT_W, Inches(0.8))
        _para(frame, s.caption, size=15, color=MUTED, italic=True, first=True,
              line_spacing=1.15)

    _notes(slide, s.notes)


def render_term(prs, s: Term, marker):
    """Plain words, then one concrete example, then the formal definition.
    Never the formal definition first."""
    slide = _blank(prs)
    _ground(slide, PAPER)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.16), GREEN)

    _chip(slide, MARGIN, Inches(0.62), "NEW TERM", fill=GOLD,
          text_color=INK, width=Inches(1.35))

    frame = _textbox(slide, MARGIN, Inches(1.15), CONTENT_W, Inches(0.9))
    _para(frame, s.term, size=32, color=GREEN, font=DISPLAY_FONT, first=True,
          line_spacing=1.05)
    _header_rule(slide, Inches(2.05))

    rows = (
        # Plain words first, one concrete example second, formal definition
        # last. Never the formal definition first.
        ("In plain words", s.plain, True),
        ("For example", s.example, False),
        ("The formal definition", s.formal, False),
    )
    inner_in = CONTENT_W.inches - 0.3
    body_pt = 19.0
    while body_pt > 14.0:
        total = sum(0.42 + _text_height(text, inner_in, body_pt,
                                        "bold" if emph else "regular")
                    for _, text, emph in rows if text)
        if total <= 4.1:
            break
        body_pt -= 1.0

    y = 2.42
    for label, text, emphasise in rows:
        if not text:
            continue
        text_h = _text_height(text, inner_in, body_pt,
                              "bold" if emphasise else "regular")
        _rect(slide, MARGIN, Inches(y + 0.05), Inches(0.07),
              Inches(text_h + 0.28), GOLD if emphasise else BORDER)
        frame = _textbox(slide, MARGIN + Inches(0.3), Inches(y),
                         CONTENT_W - Inches(0.3), Inches(text_h + 0.4))
        _para(frame, label.upper(), size=11, color=MUTED, bold=True, first=True,
              space_after=4)
        _para(frame, text, size=body_pt, color=INK, bold=emphasise,
              line_spacing=1.15)
        y += 0.42 + text_h + 0.24

    _progress(slide, marker)
    _notes(slide, s.notes)


def render_quote(prs, s: Quote, marker):
    slide = _blank(prs)
    _ground(slide, PAPER)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.16), GREEN)
    quoted = '"' + s.text + '"'
    inner_in = CONTENT_W.inches - 0.42
    quote_h = _text_height(quoted, inner_in, 25, spacing=1.2)
    source_h = _text_height(s.source, inner_in, 14) + 0.2
    _rect(slide, MARGIN, Inches(1.45), Inches(0.12),
          Inches(quote_h + source_h + 0.2), GOLD)

    frame = _textbox(slide, MARGIN, Inches(0.72), CONTENT_W, Inches(0.5))
    _para(frame, "A definition worth writing down", size=13, color=MUTED,
          bold=True, first=True)

    frame = _textbox(slide, MARGIN + Inches(0.42), Inches(1.5),
                     CONTENT_W - Inches(0.42), Inches(quote_h + source_h + 0.2))
    _para(frame, quoted, size=25, color=INK, font=DISPLAY_FONT,
          first=True, line_spacing=1.2)
    _para(frame, s.source, size=14, color=MUTED, space_before=14)

    frame = _textbox(slide, MARGIN, Inches(1.5 + quote_h + source_h + 0.55),
                     CONTENT_W, Inches(1.4))
    _para(frame, "What it tells us", size=11, color=MUTED, bold=True, first=True,
          space_after=6)
    _para(frame, s.takeaway, size=21, color=GREEN, bold=True, line_spacing=1.15)

    _progress(slide, marker)
    _notes(slide, s.notes)


def _png_size(path) -> tuple[int, int]:
    """Pixel dimensions from a PNG header, without an image library."""
    with open(path, "rb") as fh:
        head = fh.read(24)
    if head[:8] != b"\x89PNG\r\n\x1a\n":
        raise ValueError(f"{path} is not a PNG")
    return (int.from_bytes(head[16:20], "big"),
            int.from_bytes(head[20:24], "big"))


# The band a figure is drawn into, between the header rule and the credit line.
FIGURE_TOP = 1.98
FIGURE_H = 4.10
FIGURE_CAPTION_Y = 6.20


def render_figure(prs, s: Figure, marker):
    """A textbook figure, one to a slide, with a credit line under it.

    Renders identically whether or not the artwork is present. When the file
    is missing the image band carries a placeholder plate naming the figure
    and saying what it shows, so the deck built from the public repository is
    complete and the instructor can see exactly what is meant to be there.
    """
    slide = _blank(prs)
    _content_frame(slide, s.title, marker)

    band_w = CONTENT_W.inches
    path = None if FIGURES_DIR is None else Path(FIGURES_DIR) / s.filename

    if path is not None and path.is_file():
        px_w, px_h = _png_size(path)
        pad = 0.13
        scale = min((band_w - 2 * pad) / px_w, (FIGURE_H - 2 * pad) / px_h)
        w, h = px_w * scale, px_h * scale
        x = MARGIN.inches + (band_w - w) / 2.0
        y = FIGURE_TOP + (FIGURE_H - h) / 2.0
        # A white plate under the artwork: the figures are scans on white and
        # the page ground is warm paper, so an unframed image looks like a
        # patch. The plate makes the edge deliberate.
        _rect(slide, Inches(x - pad), Inches(y - pad),
              Inches(w + 2 * pad), Inches(h + 2 * pad), WHITE, line=BORDER)
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    else:
        plate_w = min(8.2, band_w)
        x = MARGIN.inches + (band_w - plate_w) / 2.0
        _rect(slide, Inches(x), Inches(FIGURE_TOP), Inches(plate_w),
              Inches(FIGURE_H), WHITE, line=BORDER)
        _rect(slide, Inches(x), Inches(FIGURE_TOP), Inches(plate_w),
              Pt(4), GOLD)

        frame = _textbox(slide, Inches(x + 0.6), Inches(FIGURE_TOP + 0.7),
                         Inches(plate_w - 1.2), Inches(FIGURE_H - 1.2))
        _para(frame, f"FIGURE {s.number}", size=15, color=MUTED, bold=True,
              first=True, space_after=14)
        _para(frame, s.shows, size=22, color=INK, font=DISPLAY_FONT,
              line_spacing=1.2, space_after=18)
        _para(frame,
              "The artwork is copyrighted and is not in the public "
              "repository. Build with assets/figures/ present to place it here.",
              size=14, color=MUTED, italic=True, line_spacing=1.15)

    frame = _textbox(slide, MARGIN, Inches(FIGURE_CAPTION_Y), CONTENT_W,
                     Inches(0.34))
    _para(frame, s.credit, size=13, color=MUTED, first=True)

    _notes(slide, s.notes)


def render_check(prs, s: Check, marker, index, total):
    """Question slide. Deliberately unlike any teaching slide: deep green
    ground, gold CHECK chip, two questions side by side."""
    slide = _blank(prs)
    _ground(slide, GREEN_DEEP)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.16), GOLD)

    _chip(slide, MARGIN, Inches(0.55), f"CHECK {index} OF {total}", fill=GOLD,
          text_color=INK, width=Inches(1.85))

    frame = _textbox(slide, MARGIN + Inches(2.1), Inches(0.57), Inches(8.0),
                     Inches(0.4))
    _para(frame, s.label, size=15, color=WHITE, bold=True, first=True)

    col_w = Emu(int((CONTENT_W - Inches(0.7)) / 2))
    inner_in = col_w.inches - 0.6
    card_top, card_h = 1.32, 5.02
    budget = card_h - 0.95  # padding, the Q label, and the gap under the stem

    # Both columns share one pair of font sizes, chosen for the worse of the
    # two questions, so a long department-style stem never overlaps its options.
    stem_pt, opt_pt = 17.0, 15.0
    for stem_try, opt_try in ((17.0, 15.0), (16.0, 14.0), (15.0, 13.5),
                              (14.0, 13.0), (13.0, 12.0), (12.0, 11.0)):
        worst = max(
            _text_height(q.stem, inner_in, stem_try, "bold")
            + sum(_text_height(f"D.  {o}", inner_in, opt_try) + 0.13
                  for o in q.options)
            for q in s.questions
        )
        stem_pt, opt_pt = stem_try, opt_try
        if worst <= budget:
            break

    tallest_stem = max(_text_height(q.stem, inner_in, stem_pt, "bold")
                       for q in s.questions)

    for i, q in enumerate(s.questions):
        x = MARGIN + (col_w + Inches(0.7)) * i
        _rect(slide, x, Inches(card_top), col_w, Inches(card_h), GREEN)

        frame = _textbox(slide, x + Inches(0.3), Inches(card_top + 0.27),
                         col_w - Inches(0.6), Inches(0.25))
        _para(frame, f"Q{i + 1}", size=12, color=GOLD, bold=True, first=True)

        stem_top = card_top + 0.58
        stem_h = _text_height(q.stem, inner_in, stem_pt, "bold")
        frame = _textbox(slide, x + Inches(0.3), Inches(stem_top),
                         col_w - Inches(0.6), Inches(stem_h + 0.1))
        _para(frame, q.stem, size=stem_pt, color=WHITE, bold=True, first=True,
              line_spacing=1.15)

        # Options start on a shared baseline so the two columns read as a pair
        # even when one stem is much longer than the other.
        frame = _textbox(slide, x + Inches(0.3),
                         Inches(stem_top + tallest_stem + 0.32),
                         col_w - Inches(0.6), Inches(3.2))
        for j, opt in enumerate(q.options):
            _para(frame, f"{'ABCD'[j]}.  {opt}", size=opt_pt, color=WHITE,
                  first=(j == 0), space_after=9, line_spacing=1.12)

    _progress(slide, marker, on_dark=True)
    _notes(slide, s.notes or (
        "Read Q1 aloud, then Q2. Give them sixty seconds and no discussion.",
        "Ask for a show of hands on each option before you advance.",
        "This is for reading the room, not for marks. Say so if anyone tenses up.",
    ))


def render_reveal(prs, s: Check, marker, index, total):
    """Answer slide. Paper ground so the room feels the release, with the
    answers in the monospace face and one line of reason each."""
    slide = _blank(prs)
    _ground(slide, PAPER)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.16), GOLD)

    _chip(slide, MARGIN, Inches(0.55), "ANSWERS", fill=GREEN,
          text_color=WHITE, width=Inches(1.35))

    frame = _textbox(slide, MARGIN + Inches(1.6), Inches(0.57), Inches(8.5),
                     Inches(0.4))
    _para(frame, s.label, size=15, color=GREEN, bold=True, first=True)

    inner_in = CONTENT_W.inches - 0.8
    available = 5.1  # from y = 1.32 down to the progress marker

    # Shrink together until both cards fit. The answer stays in the monospace
    # face and stays the largest thing on the slide.
    stem_pt, ans_pt, reason_pt = 16.0, 20.0, 16.0
    for stem_try, ans_try, reason_try in ((16.0, 20.0, 16.0), (15.0, 18.0, 15.0),
                                          (14.0, 17.0, 14.0), (13.0, 15.0, 13.0)):
        heights = [
            0.95
            + _text_height(f"Q1   {q.stem}", inner_in, stem_try)
            + _text_height(f"{q.answer}.  {q.options['ABCD'.index(q.answer)]}",
                           inner_in, ans_try, "mono")
            + _text_height(q.reason, inner_in, reason_try)
            for q in s.questions
        ]
        stem_pt, ans_pt, reason_pt = stem_try, ans_try, reason_try
        if sum(heights) + 0.27 <= available:
            break

    y = 1.32
    for i, q in enumerate(s.questions):
        stem = f"Q{i + 1}   {q.stem}"
        answer = f"{q.answer}.  {q.options['ABCD'.index(q.answer)]}"
        stem_h = _text_height(stem, inner_in, stem_pt)
        ans_h = _text_height(answer, inner_in, ans_pt, "mono")
        card_h = 0.95 + stem_h + ans_h + _text_height(q.reason, inner_in, reason_pt)

        _rect(slide, MARGIN, Inches(y), CONTENT_W, Inches(card_h), WHITE,
              line=BORDER)
        _rect(slide, MARGIN, Inches(y), Inches(0.1), Inches(card_h), GOLD)

        frame = _textbox(slide, MARGIN + Inches(0.4), Inches(y + 0.26),
                         CONTENT_W - Inches(0.8), Inches(stem_h + 0.1))
        _para(frame, stem, size=stem_pt, color=MUTED, first=True,
              line_spacing=1.15)

        frame = _textbox(slide, MARGIN + Inches(0.4),
                         Inches(y + 0.26 + stem_h + 0.28),
                         CONTENT_W - Inches(0.8), Inches(ans_h + 1.0))
        _para(frame, answer, size=ans_pt, color=GREEN, font=MONO_FONT,
              bold=True, first=True, line_spacing=1.15)
        _para(frame, q.reason, size=reason_pt, color=INK, space_before=7,
              line_spacing=1.15)

        y += card_h + 0.27

    _progress(slide, marker)
    _notes(slide, s.notes or (
        "Reveal one at a time. Say the reason out loud, not just the letter.",
        "If more than a third of the room missed one, step back and reteach it.",
    ))


def render_recap(prs, s: Recap, marker):
    slide = _blank(prs)
    _ground(slide, PAPER)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.16), GREEN)
    _rect(slide, SLIDE_W - Inches(0.28), 0, Inches(0.28), SLIDE_H, GOLD)

    frame = _textbox(slide, MARGIN, Inches(0.72), CONTENT_W - Inches(0.5),
                     Inches(0.9))
    _para(frame, s.title, size=32, color=GREEN, font=DISPLAY_FONT, first=True)
    _header_rule(slide, Inches(1.72))

    col_w = Emu(int((CONTENT_W - Inches(1.1)) / 2))
    inner_in = col_w.inches
    available = 4.15  # from y = 2.15 down to the progress marker

    # Recap lists vary in length between sections, so size the type to the
    # list and split the two columns by height rather than by item count.
    size, gap = 18.0, 0.17
    columns = ((), ())
    for try_size in (18.0, 17.0, 16.0, 15.0, 14.0, 13.0):
        heights = [_text_height("- " + item, inner_in, try_size) + gap
                   for item in s.items]
        target = sum(heights) / 2.0
        split, running = len(s.items), 0.0
        for i, h in enumerate(heights):
            if running + h / 2 > target:
                split = i
                break
            running += h
        left, right = heights[:split], heights[split:]
        size = try_size
        columns = (s.items[:split], s.items[split:])
        if max(sum(left), sum(right)) <= available:
            break

    for c, items in enumerate(columns):
        if not items:
            continue
        x = MARGIN + (col_w + Inches(0.6)) * c
        frame = _textbox(slide, x, Inches(2.15), col_w, Inches(available))
        for i, item in enumerate(items):
            _para(frame, "\u25aa  " + item, size=size, color=INK,
                  first=(i == 0), space_after=int(gap * 72), line_spacing=1.12,
                  hanging=Inches(0.26))

    _progress(slide, marker)
    _notes(slide, s.notes)


def render_closing(prs, s: Closing):
    slide = _blank(prs)
    _ground(slide, GREEN)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.22), GOLD)

    frame = _textbox(slide, MARGIN, Inches(1.5), CONTENT_W, Inches(1.1))
    _para(frame, s.title, size=38, color=WHITE, font=DISPLAY_FONT, first=True,
          line_spacing=1.05)
    _rect(slide, MARGIN, Inches(2.9), Inches(1.6), Pt(3), GOLD)

    frame = _textbox(slide, MARGIN, Inches(3.4), CONTENT_W, Inches(2.4))
    for i, line in enumerate(s.lines):
        _para(frame, line, size=20, color=WHITE, first=(i == 0), space_after=12,
              line_spacing=1.15)

    if s.accent:
        frame = _textbox(slide, MARGIN, SLIDE_H - Inches(1.5), CONTENT_W,
                         Inches(0.7))
        _para(frame, s.accent, size=18, color=GOLD, bold=True, first=True)

    _notes(slide, s.notes)


# --------------------------------------------------------------------------
# Build
# --------------------------------------------------------------------------

MAX_BODY_LINES = 6

# Nothing may be drawn below this line: the progress marker sits at 6.78in.
SAFE_BOTTOM = 6.50


def _content_bottom(s: Content) -> float:
    """Where render_content will finish drawing, in inches."""
    inner = CONTENT_W.inches
    y = 2.15
    if s.lines:
        y += sum(_text_height(line, inner, 21, spacing=1.18) + 0.20
                 for line in s.lines) + 0.34
    if s.accent:
        y += _text_height(s.accent, inner - 0.32, 22, "bold") + 0.16 + 0.26
    if s.caption:
        y += _text_height(s.caption, inner, 15)
    return y


def _term_bottom(s: Term) -> float:
    """Where render_term will finish drawing, in inches, at the smallest
    body size it is allowed to fall back to."""
    inner = CONTENT_W.inches - 0.3
    rows = ((s.plain, True), (s.example, False), (s.formal, False))
    for pt in (19.0, 18.0, 17.0, 16.0, 15.0, 14.0):
        total = sum(0.42 + _text_height(text, inner, pt,
                                        "bold" if emph else "regular") + 0.24
                    for text, emph in rows if text)
        if total <= 4.1 or pt == 14.0:
            return 2.42 + total
    return 2.42


def _quote_bottom(s: Quote) -> float:
    inner = CONTENT_W.inches - 0.42
    quote_h = _text_height('"' + s.text + '"', inner, 25, spacing=1.2)
    source_h = _text_height(s.source, inner, 14) + 0.2
    top = 1.5 + quote_h + source_h + 0.55
    return top + 0.3 + _text_height(s.takeaway, CONTENT_W.inches, 21, "bold")


def _validate(chapter: Chapter) -> list[str]:
    """Enforce the teaching brief's hard rules at build time, so a future
    chapter cannot quietly regress them."""
    problems: list[str] = []

    def check_text(where: str, text: str):
        if "—" in text or "–" in text:
            problems.append(f"{where}: dash character not allowed, use a plain dash")

    def walk(where, slide):
        for value in vars(slide).values():
            if isinstance(value, str):
                check_text(where, value)
            elif isinstance(value, (tuple, list)):
                for item in value:
                    if isinstance(item, str):
                        check_text(where, item)

    for section in chapter.sections:
        for i, slide in enumerate(section.slides):
            where = f"section {section.number} slide {i + 1}"
            walk(where, slide)
            if isinstance(slide, Content) and len(slide.lines) > MAX_BODY_LINES:
                problems.append(
                    f"{where}: {len(slide.lines)} body lines, limit is {MAX_BODY_LINES}"
                )
            if isinstance(slide, Figure):
                if not slide.number or "." not in slide.number:
                    problems.append(f"{where}: a figure needs a book figure number")
                if not slide.shows:
                    problems.append(
                        f"{where}: figure {slide.number} needs a 'shows' line, "
                        "which is what the placeholder prints without the artwork"
                    )
            if isinstance(slide, Check):
                if len(slide.questions) != 2:
                    problems.append(f"{where}: a check must hold exactly 2 questions")
                for q in slide.questions:
                    walk(where, q)
                    if len(q.options) != 4:
                        problems.append(f"{where}: every question needs 4 options")
                    if q.answer not in "ABCD":
                        problems.append(f"{where}: answer must be A, B, C or D")
            bottom = None
            if isinstance(slide, Content):
                bottom = _content_bottom(slide)
            elif isinstance(slide, Term):
                bottom = _term_bottom(slide)
            elif isinstance(slide, Quote):
                bottom = _quote_bottom(slide)
            if bottom is not None and bottom > SAFE_BOTTOM:
                problems.append(
                    f"{where}: content runs to {bottom:.2f}in, past the "
                    f"{SAFE_BOTTOM}in safe bottom. Split the slide or shorten it."
                )
            if not isinstance(slide, Check) and not slide.notes:
                problems.append(f"{where}: content slides need speaker cues")
            if slide.notes and len(slide.notes) > 3:
                problems.append(f"{where}: keep speaker cues to two or three lines")
        walk(f"section {section.number} recap", section.recap)

    # Answer positions must be spread across A/B/C/D. A skewed key lets a
    # student score well by always picking one letter, which destroys the
    # only thing the checks are for: reading the room.
    letters = [q.answer
               for section in chapter.sections
               for slide in section.slides if isinstance(slide, Check)
               for q in slide.questions]
    if letters:
        for L in "ABCD":
            share = letters.count(L) / len(letters)
            if share > 0.35:
                problems.append(
                    f"answer key: {L} holds {share:.0%} of {len(letters)} items, "
                    "limit is 35%. Reorder options to spread the key."
                )
            elif share < 0.15:
                problems.append(
                    f"answer key: {L} holds only {share:.0%} of {len(letters)} "
                    "items, floor is 15%. Reorder options to spread the key."
                )
        runs = [i for i in range(len(letters) - 2)
                if letters[i] == letters[i + 1] == letters[i + 2]]
        if runs:
            problems.append(
                f"answer key: {len(runs)} run(s) of three identical answers in a "
                "row, starting at item(s) "
                + ", ".join(str(i + 1) for i in runs) + "."
            )
    return problems


def build(chapter: Chapter, out_path, *, display_font: str | None = None,
          figures_dir=None):
    """Render the chapter to a .pptx. Deterministic: same input, same deck.

    ``figures_dir`` is where render_figure looks for the textbook artwork. Pass
    None, or a path that does not exist, and every figure slide renders as a
    placeholder instead. The slide count and the progress markers are the same
    either way.
    """
    global DISPLAY_FONT, FIGURES_DIR
    if display_font:
        DISPLAY_FONT = display_font
    FIGURES_DIR = Path(figures_dir) if figures_dir else None

    problems = _validate(chapter)
    if problems:
        raise SystemExit(
            "Teaching design rules violated:\n  " + "\n  ".join(problems)
        )

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Fixed document properties keep the build byte-for-byte reproducible, so
    # rebuilding an unchanged chapter produces no diff.
    core = prs.core_properties
    core.title = f"{chapter.code} {chapter.chapter} - {chapter.title}"
    core.author = chapter.presenter
    core.last_modified_by = chapter.presenter
    core.subject = chapter.course
    core.revision = 1
    core.created = EPOCH
    core.modified = EPOCH

    total_sections = len(chapter.sections)
    checks: list[tuple[int, Check]] = []
    for section in chapter.sections:
        for slide in section.slides:
            if isinstance(slide, Check):
                checks.append((section.number, slide))
    total_checks = len(checks)

    render_title(prs, TitleSlide(
        course=f"{chapter.code}  {chapter.course}",
        chapter=chapter.chapter,
        title=chapter.title,
        subtitle=chapter.subtitle,
        presenter=chapter.presenter,
        notes=(
            "Greet the room, then say what today buys them: after this session "
            "they can read a chart argument and say what it does and does not claim.",
            "Point at the roadmap slide next and promise a clean stopping point "
            "every twenty five minutes.",
        ),
    ))

    # The booklet lists seven objectives, which is more than one slide should
    # carry, so they are split rather than crammed.
    objective_chunks = [chapter.objectives[i:i + 4]
                        for i in range(0, len(chapter.objectives), 4)]
    for n, chunk in enumerate(objective_chunks):
        render_content(prs, Content(
            title="What you will be able to do" if n == 0 else "What you will be able to do, continued",
            lines=chunk,
            notes=(
                "Read these out once, slowly. These are the exact learning "
                "objectives in the course booklet for this week.",
                "Tell them every one of these is examinable.",
            ),
        ), marker=f"{chapter.chapter} - Learning objectives")

    render_content(prs, Content(
        title="How today is laid out",
        lines=chapter.roadmap,
        accent="Six parts, a check every few terms, and a clean stop at every boundary.",
        notes=(
            "Say the shape of the session out loud: six parts, roughly twenty "
            "five minutes each.",
            "Tell them the checks are for reading the room, not for marks.",
        ),
    ), marker=f"{chapter.chapter} - Roadmap")

    check_index = 0
    for section in chapter.sections:
        render_section_open(prs, SectionOpen(
            number=section.number,
            total=total_sections,
            title=section.title,
            covers=section.covers,
            minutes=section.minutes,
            notes=(
                f"Say we are starting Part {section.number} of {total_sections}.",
                "Read the three covers lines, then move. Do not linger here.",
            ),
        ))

        position = 0
        body = [s for s in section.slides if not isinstance(s, Check)]
        for slide in section.slides:
            if isinstance(slide, Check):
                check_index += 1
                marker = (f"Part {section.number} of {total_sections} - "
                          f"{section.short}  |  Check {check_index}")
                render_check(prs, slide, marker, check_index, total_checks)
                render_reveal(prs, slide, marker, check_index, total_checks)
                continue
            position += 1
            marker = (f"Part {section.number} of {total_sections} - "
                      f"{section.short}  |  {position} of {len(body)}")
            if isinstance(slide, Term):
                render_term(prs, slide, marker)
            elif isinstance(slide, Quote):
                render_quote(prs, slide, marker)
            elif isinstance(slide, Content):
                render_content(prs, slide, marker)
            elif isinstance(slide, Figure):
                render_figure(prs, slide, marker)
            else:
                raise TypeError(f"unhandled slide type {type(slide).__name__}")

        render_recap(prs, section.recap,
                     marker=(f"Part {section.number} of {total_sections} - "
                             f"{section.short}  |  recap"))

    for slide in chapter.closing:
        if isinstance(slide, Closing):
            render_closing(prs, slide)
        elif isinstance(slide, Content):
            render_content(prs, slide, marker=f"{chapter.chapter} - Wrap up")
        else:
            raise TypeError(f"unhandled closing slide {type(slide).__name__}")

    prs.save(str(out_path))
    _repack_deterministically(out_path)
    return len(prs.slides._sldIdLst), total_checks


def figure_status(chapter: Chapter, figures_dir=None):
    """Every figure the chapter declares, and whether its artwork is present."""
    root = Path(figures_dir) if figures_dir else None
    out = []
    for section in chapter.sections:
        for slide in section.slides:
            if isinstance(slide, Figure):
                have = root is not None and (root / slide.filename).is_file()
                out.append((slide.number, slide.filename, have))
    return out


def _repack_deterministically(path) -> None:
    """Rewrite the .pptx zip with a stable entry order and a fixed timestamp.

    python-pptx does not guarantee a stable order for the package parts, so
    two builds of identical content can differ byte for byte. That produces
    noisy diffs on a generated file that is committed. Repacking fixes it.
    """
    path = str(path)
    with zipfile.ZipFile(path) as src:
        entries = {name: src.read(name) for name in src.namelist()}

    # [Content_Types].xml conventionally comes first; the rest sort by name.
    first = "[Content_Types].xml"
    names = ([first] if first in entries else []) + sorted(
        n for n in entries if n != first)

    tmp = path + ".tmp"
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for name in names:
            info = zipfile.ZipInfo(name, date_time=EPOCH.timetuple()[:6])
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            dst.writestr(info, entries[name])
    os.replace(tmp, path)


def iter_checks(chapter: Chapter):
    """Every check in reading order, for the instructor's answer sheet."""
    index = 0
    for section in chapter.sections:
        for slide in section.slides:
            if isinstance(slide, Check):
                index += 1
                yield index, section, slide
